from docx import Document 
from pdfminer.high_level import extract_text
from pptx import Presentation

def read_doc(stream_address):
    doc = Document(stream_address)
    return '\n'.join(paragraph.text for paragraph in doc.paragraphs)

def read_pdf(stream_address):
    return extract_text(stream_address) 

def read_ppt(stream_address):
    pres = Presentation(stream_address)
    return '\n'.join([shape for slide in pres.slides for shape in slide.shapes if hasattr(shape, 'text')])

def read_txt(stream_address):
    with open(stream_address, 'r') as file:
        text = file.read()
        return text