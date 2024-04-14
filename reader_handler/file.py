from dataclasses import dataclass
from enum import Enum
from abc import abstractmethod
from docx import Document 
from pdfminer.high_level import extract_text
from pptx import Presentation


from . import file_validator

class FileExceptions(Enum):
    """ Stream exceptions
    """
    INVALID_EXTENSION = 'Extension is not valid.'
    MISSING_FILE = 'File is missing.'

class InvalidFile(Exception):
    pass

@dataclass(frozen=True)
class File:
    file_address: str

    def __post_init__(self):
        self._validate_file()

    def _validate_file(self) -> None:
        validation_functions = (
            (FileExceptions.INVALID_EXTENSION, lambda file_address: file_validator.validate_extension(file_address)),
            (FileExceptions.MISSING_FILE, lambda file_address: file_validator.validate_availability(file_address))
            )
        for exception, validation_func in validation_functions:
            error_flag, error_message = validation_func(self.stream_address)
            if error_flag:
                raise InvalidFile('Exception: {}. ErrorMessage: {}'.format(
                    exception.value, 
                    error_message))
            
    def gen_read_binary(self) -> str:
        return None
    
    @abstractmethod 
    def read_file(self) -> str:
        raise NotImplementedError
            
class PDFDocument(File):
    def read_file(self) -> str:
        return extract_text(self.file_address)

class WordDocument(File):
    def read_file(self) -> str:
        doc = Document(self.file_address)
        return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
    
class PresentationDocument(File):
    def read_file(self) -> str:
        pres = Presentation(self.file_address)
        return '\n'.join([shape for slide in pres.slides for shape in slide.shapes if hasattr(shape, 'text')])
    
class FlatDocument(File):
    def read_file(self) -> str:
        with open(self.file_address, 'r') as file:
            text = file.read()
            return text 




